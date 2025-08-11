import os
import json
import tkinter as tk
from tkinter import filedialog, ttk, messagebox, scrolledtext
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import pygments
from pygments import lexers, formatters
from bs4 import BeautifulSoup
import markdown2
import re
import datetime


# Mobile-optimized themes with better contrast
THEMES = {
    "Light": {"bg_color": "#FFFFFF", "text_color": "#000000", "accent": "#007ACC"},
    "Dark": {"bg_color": "#1E1E1E", "text_color": "#FFFFFF", "accent": "#4FC3F7"},
    "Blue": {"bg_color": "#0D47A1", "text_color": "#FFFFFF", "accent": "#81D4FA"},
    "Green": {"bg_color": "#1B5E20", "text_color": "#FFFFFF", "accent": "#81C784"},
    "Purple": {"bg_color": "#4A148C", "text_color": "#FFFFFF", "accent": "#CE93D8"}
}

# Programming language extensions mapping
LANGUAGE_EXTENSIONS = {
    'python': '.py', 'javascript': '.js', 'java': '.java', 'cpp': '.cpp', 'c++': '.cpp',
    'c': '.c', 'csharp': '.cs', 'c#': '.cs', 'php': '.php', 'ruby': '.rb', 'go': '.go',
    'rust': '.rs', 'swift': '.swift', 'kotlin': '.kt', 'scala': '.scala', 'r': '.r',
    'matlab': '.m', 'sql': '.sql', 'html': '.html', 'css': '.css', 'typescript': '.ts',
    'bash': '.sh', 'shell': '.sh', 'powershell': '.ps1', 'dockerfile': 'Dockerfile',
    'yaml': '.yml', 'yml': '.yml', 'json': '.json', 'xml': '.xml', 'markdown': '.md',
    'tex': '.tex', 'latex': '.tex', 'vim': '.vim', 'lua': '.lua', 'perl': '.pl',
    'dart': '.dart', 'haskell': '.hs', 'clojure': '.clj', 'elixir': '.ex', 'erlang': '.erl',
    'f#': '.fs', 'groovy': '.groovy', 'julia': '.jl', 'objc': '.m', 'pascal': '.pas',
    'fortran': '.f90', 'assembly': '.asm', 'vhdl': '.vhd', 'verilog': '.v'
}


class FinalPPTCreatorApp:
    def __init__(self, root):
        self.root = root
        self.root.title("📱 Final Mobile PPT Generator - Multi JSON Support + Code Files")
        self.root.geometry("700x650")
        self.all_topics_data = []
        self.loaded_files = []
        
        # Setup UI
        self.setup_ui()
        
    def setup_ui(self):
        # Main frame with scrollbar
        main_canvas = tk.Canvas(self.root)
        scrollbar = ttk.Scrollbar(self.root, orient="vertical", command=main_canvas.yview)
        scrollable_frame = ttk.Frame(main_canvas)

        scrollable_frame.bind(
            "<Configure>",
            lambda e: main_canvas.configure(scrollregion=main_canvas.bbox("all"))
        )

        main_canvas.create_window((0, 0), window=scrollable_frame, anchor="nw")
        main_canvas.configure(yscrollcommand=scrollbar.set)
        
        main_frame = ttk.Frame(scrollable_frame, padding="20")
        main_frame.pack(fill=tk.BOTH, expand=True)
        
        # Title
        title_label = ttk.Label(main_frame, text="📱 Final Mobile PPT Generator", 
                               font=("Arial", 18, "bold"))
        title_label.pack(pady=(0, 5))
        
        subtitle_label = ttk.Label(main_frame, text="Multi-JSON Support -  Organized Folders -  Mobile Optimized -  Auto Code Files", 
                                  font=("Arial", 10), foreground="gray")
        subtitle_label.pack(pady=(0, 20))
        
        # Feature highlight
        feature_frame = ttk.Frame(main_frame)
        feature_frame.pack(fill=tk.X, pady=(0, 15))
        
        feature_label = ttk.Label(feature_frame, text="🆕 NEW: Automatically creates separate code files for slides with 'slide_type': 'code'", 
                                 font=("Arial", 9, "bold"), foreground="green")
        feature_label.pack()
        
        # JSON Input Section
        json_frame = ttk.LabelFrame(main_frame, text="📄 JSON Input Management", padding="15")
        json_frame.pack(fill=tk.X, pady=(0, 20))
        
        # Buttons frame
        buttons_frame = ttk.Frame(json_frame)
        buttons_frame.pack(fill=tk.X, pady=(0, 15))
        
        # Load multiple JSON files button
        load_multiple_btn = tk.Button(buttons_frame, text="📂 Load Multiple JSON Files", 
                                     font=("Arial", 11, "bold"), bg="#28a745", fg="white",
                                     command=self.load_multiple_json_files, cursor="hand2")
        load_multiple_btn.pack(side=tk.LEFT, padx=(0, 10))
        
        # Clear files button
        clear_btn = tk.Button(buttons_frame, text="🗑️ Clear All", 
                             font=("Arial", 11), bg="#dc3545", fg="white",
                             command=self.clear_all_files, cursor="hand2")
        clear_btn.pack(side=tk.RIGHT)
        
        # JSON Text Input Area
        text_input_frame = ttk.LabelFrame(json_frame, text="✏️ Paste JSON Content", padding="10")
        text_input_frame.pack(fill=tk.BOTH, expand=True, pady=(10, 0))
        
        # Instructions
        instruction_label = ttk.Label(text_input_frame, 
                                    text="Paste your JSON array content below and click 'Add from Text':",
                                    font=("Arial", 10))
        instruction_label.pack(anchor=tk.W, pady=(0, 5))
        
        # Text area for JSON input
        self.json_text_area = scrolledtext.ScrolledText(text_input_frame, 
                                                       height=8, 
                                                       width=60,
                                                       font=("Consolas", 10),
                                                       wrap=tk.WORD)
        self.json_text_area.pack(fill=tk.BOTH, expand=True, pady=(0, 10))
        
        # Add placeholder text with code example
        placeholder_text = '''[
  {
    "topic": "Python Basics",
    "slides": [
      {
        "title": "Hello World Example",
        "content": "print('Hello, World!')",
        "slide_type": "code",
        "language": "python"
      },
      {
        "title": "Variables in Python",
        "content": "Variables are containers for storing data values.",
        "slide_type": "text"
      }
    ]
  }
]'''
        self.json_text_area.insert(tk.END, placeholder_text)
        self.json_text_area.bind("<FocusIn>", self.clear_placeholder)
        
        # Button to add from text
        text_buttons_frame = ttk.Frame(text_input_frame)
        text_buttons_frame.pack(fill=tk.X)
        
        add_text_btn = tk.Button(text_buttons_frame, text="➕ Add from Text", 
                                font=("Arial", 11, "bold"), bg="#17a2b8", fg="white",
                                command=self.add_from_text, cursor="hand2")
        add_text_btn.pack(side=tk.LEFT, padx=(0, 10))
        
        clear_text_btn = tk.Button(text_buttons_frame, text="🧹 Clear Text", 
                                  font=("Arial", 10), bg="#6c757d", fg="white",
                                  command=self.clear_text_area, cursor="hand2")
        clear_text_btn.pack(side=tk.LEFT)
        
        # Files list frame
        files_frame = ttk.Frame(json_frame)
        files_frame.pack(fill=tk.BOTH, expand=True, pady=(15, 0))
        
        ttk.Label(files_frame, text="Loaded Sources:", font=("Arial", 10, "bold")).pack(anchor=tk.W, pady=(0, 5))
        
        # Files listbox with scrollbar
        files_list_frame = ttk.Frame(files_frame)
        files_list_frame.pack(fill=tk.BOTH, expand=True)
        
        self.files_listbox = tk.Listbox(files_list_frame, height=4, font=("Arial", 9))
        files_scroll = ttk.Scrollbar(files_list_frame, orient="vertical", command=self.files_listbox.yview)
        self.files_listbox.configure(yscrollcommand=files_scroll.set)
        self.files_listbox.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)
        files_scroll.pack(side=tk.RIGHT, fill=tk.Y)
        
        # Configuration Section
        config_frame = ttk.LabelFrame(main_frame, text="⚙️ Configuration", padding="15")
        config_frame.pack(fill=tk.X, pady=(0, 20))
        
        # Theme selection
        theme_frame = ttk.Frame(config_frame)
        theme_frame.pack(fill=tk.X, pady=(0, 15))
        
        ttk.Label(theme_frame, text="🎨 Select Theme:", font=("Arial", 10, "bold")).pack(anchor=tk.W)
        self.theme_var = tk.StringVar(value="Light")
        self.theme_dropdown = ttk.Combobox(theme_frame, textvariable=self.theme_var, 
                                          values=list(THEMES.keys()), state="readonly", width=20)
        self.theme_dropdown.pack(anchor=tk.W, pady=(5, 0))
        
        # Code file generation option
        code_frame = ttk.Frame(config_frame)
        code_frame.pack(fill=tk.X, pady=(0, 15))
        
        ttk.Label(code_frame, text="💻 Code File Generation:", font=("Arial", 10, "bold")).pack(anchor=tk.W)
        self.generate_code_files_var = tk.BooleanVar(value=True)
        code_check = ttk.Checkbutton(code_frame, text="Generate separate code files for 'slide_type': 'code'", 
                                    variable=self.generate_code_files_var)
        code_check.pack(anchor=tk.W, pady=(5, 0))
        
        # Starting folder number
        folder_frame = ttk.Frame(config_frame)
        folder_frame.pack(fill=tk.X, pady=(0, 15))
        
        ttk.Label(folder_frame, text="📁 Starting Folder Number:", font=("Arial", 10, "bold")).pack(anchor=tk.W)
        self.folder_number_var = tk.StringVar(value="1")
        folder_number_entry = ttk.Entry(folder_frame, textvariable=self.folder_number_var, width=10)
        folder_number_entry.pack(anchor=tk.W, pady=(5, 0))
        
        # Folder naming option
        naming_frame = ttk.Frame(config_frame)
        naming_frame.pack(fill=tk.X, pady=(0, 15))
        
        ttk.Label(naming_frame, text="📋 Folder Naming:", font=("Arial", 10, "bold")).pack(anchor=tk.W)
        self.naming_var = tk.StringVar(value="numbered")
        naming_radio_frame = ttk.Frame(naming_frame)
        naming_radio_frame.pack(anchor=tk.W, pady=(5, 0))
        
        ttk.Radiobutton(naming_radio_frame, text="01. Topic Name", 
                       variable=self.naming_var, value="numbered").pack(side=tk.LEFT, padx=(0, 15))
        ttk.Radiobutton(naming_radio_frame, text="Topic Name Only", 
                       variable=self.naming_var, value="name_only").pack(side=tk.LEFT)
        
        # Output directory selection
        output_frame = ttk.Frame(config_frame)
        output_frame.pack(fill=tk.X, pady=(0, 10))
        
        ttk.Label(output_frame, text="📂 Output Directory:", font=("Arial", 10, "bold")).pack(anchor=tk.W)
        
        dir_select_frame = ttk.Frame(output_frame)
        dir_select_frame.pack(fill=tk.X, pady=(5, 0))
        
        self.output_dir_var = tk.StringVar(value=os.getcwd())
        output_dir_entry = ttk.Entry(dir_select_frame, textvariable=self.output_dir_var, state="readonly")
        output_dir_entry.pack(side=tk.LEFT, fill=tk.X, expand=True, padx=(0, 10))
        
        browse_btn = tk.Button(dir_select_frame, text="Browse", 
                              font=("Arial", 9), command=self.browse_output_directory)
        browse_btn.pack(side=tk.RIGHT)
        
        # Topics Preview Section
        preview_frame = ttk.LabelFrame(main_frame, text="📋 All Topics Preview", padding="15")
        preview_frame.pack(fill=tk.X, pady=(0, 20))
        
        # Topics info frame
        topics_info_frame = ttk.Frame(preview_frame)
        topics_info_frame.pack(fill=tk.X, pady=(0, 10))
        
        self.topics_info_label = ttk.Label(topics_info_frame, text="No topics loaded", 
                                          font=("Arial", 10), foreground="gray")
        self.topics_info_label.pack(anchor=tk.W)
        
        # Topics listbox
        topics_list_frame = ttk.Frame(preview_frame)
        topics_list_frame.pack(fill=tk.BOTH, expand=True)
        
        self.topics_listbox = tk.Listbox(topics_list_frame, height=6, font=("Arial", 9))
        topics_scroll = ttk.Scrollbar(topics_list_frame, orient="vertical", command=self.topics_listbox.yview)
        self.topics_listbox.configure(yscrollcommand=topics_scroll.set)
        self.topics_listbox.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)
        topics_scroll.pack(side=tk.RIGHT, fill=tk.Y)
        
        # Generate button
        generate_btn = tk.Button(main_frame, text="🚀 Generate All PPTs with Organized Folders + Code Files", 
                               font=("Arial", 14, "bold"), bg="#007ACC", fg="white",
                               command=self.generate_all_presentations, cursor="hand2",
                               height=2)
        generate_btn.pack(pady=20, fill=tk.X)
        
        # Status label
        self.status_label = tk.Label(main_frame, text="Ready to create organized presentations with automatic code file generation", 
                                   font=("Arial", 10), fg="gray")
        self.status_label.pack(pady=(10, 0))

        # Pack canvas and scrollbar
        main_canvas.pack(side="left", fill="both", expand=True)
        scrollbar.pack(side="right", fill="y")

    def clear_placeholder(self, event):
        """Clear placeholder text when text area gains focus"""
        current_text = self.json_text_area.get("1.0", tk.END).strip()
        if current_text.startswith('['):
            # Check if it's still the placeholder
            try:
                data = json.loads(current_text)
                if (isinstance(data, list) and len(data) == 1 and 
                    data[0].get("topic") == "Python Basics"):
                    self.json_text_area.delete("1.0", tk.END)
            except:
                pass

    def clear_text_area(self):
        """Clear the text area"""
        self.json_text_area.delete("1.0", tk.END)

    def add_from_text(self):
        """Add topics from the text area content"""
        json_text = self.json_text_area.get("1.0", tk.END).strip()
        
        if not json_text:
            messagebox.showwarning("Warning", "Please paste JSON content in the text area")
            return
        
        try:
            # Parse JSON
            data = json.loads(json_text)
            
            # Validate JSON structure
            if not self.validate_json_structure(data, "Text Input"):
                return
            
            # Add source info to each topic
            for topic in data:
                topic["_source_file"] = "Text Input"
            
            # Add to all topics data
            self.all_topics_data.extend(data)
            
            # Add to loaded files list for tracking
            source_name = f"Text Input ({len(data)} topics)"
            if source_name not in self.loaded_files:
                self.loaded_files.append(source_name)
            
            # Update UI
            self.update_ui_after_load()
            
            # Clear text area after successful load
            self.json_text_area.delete("1.0", tk.END)
            
            self.status_label.config(text=f"✅ Added {len(data)} topics from text input", fg="green")
            
        except json.JSONDecodeError as e:
            messagebox.showerror("JSON Error", f"Invalid JSON format:\n{str(e)}")
            self.status_label.config(text="❌ Invalid JSON format", fg="red")
        except Exception as e:
            messagebox.showerror("Error", f"Error processing text input:\n{str(e)}")
            self.status_label.config(text="❌ Error processing input", fg="red")

    def load_multiple_json_files(self):
        """Load multiple JSON files at once"""
        file_paths = filedialog.askopenfilenames(
            filetypes=[("JSON Files", "*.json")], 
            title="Select Multiple JSON Files"
        )
        if not file_paths:
            return

        successful_loads = 0
        failed_loads = []
        
        for file_path in file_paths:
            try:
                filename = os.path.basename(file_path)
                if filename not in [f.split(' (')[0] for f in self.loaded_files if not f.startswith("Text Input")]:
                    success = self.load_json_data(file_path)
                    if success:
                        successful_loads += 1
                        # Count topics from this file for display
                        topics_count = len([t for t in self.all_topics_data if t.get("_source_file") == filename])
                        self.loaded_files.append(f"{filename} ({topics_count} topics)")
                    else:
                        failed_loads.append(filename)
                else:
                    self.status_label.config(text=f"⚠️ {filename} already loaded", fg="orange")
            except Exception as e:
                failed_loads.append(f"{os.path.basename(file_path)}: {str(e)}")
        
        self.update_ui_after_load()
        
        # Status message
        if successful_loads > 0:
            msg = f"✅ Loaded {successful_loads} files"
            if failed_loads:
                msg += f", {len(failed_loads)} failed"
            self.status_label.config(text=msg, fg="green" if not failed_loads else "orange")
        elif failed_loads:
            self.status_label.config(text=f"❌ Failed to load {len(failed_loads)} files", fg="red")
            messagebox.showerror("Load Error", f"Failed files:\n" + "\n".join(failed_loads[:3]))

    def load_json_data(self, file_path):
        """Load and validate JSON data from a file"""
        try:
            with open(file_path, "r", encoding="utf-8") as file:
                data = json.load(file)
            
            # Validate JSON structure
            if not self.validate_json_structure(data, file_path):
                return False
            
            # Add file info to each topic for tracking
            for topic in data:
                topic["_source_file"] = os.path.basename(file_path)
            
            # Add to all topics data
            self.all_topics_data.extend(data)
            return True
            
        except Exception as e:
            messagebox.showerror("JSON Error", f"Error in {os.path.basename(file_path)}:\n{str(e)}")
            return False

    def validate_json_structure(self, data, source):
        """Validate the JSON structure"""
        source_name = source if isinstance(source, str) and not os.path.exists(source) else os.path.basename(source)
        
        if not isinstance(data, list):
            messagebox.showerror("Invalid JSON", f"{source_name}: JSON should be a list of topics")
            return False
        
        for i, topic in enumerate(data):
            if not isinstance(topic, dict):
                messagebox.showerror("Invalid JSON", f"{source_name}: Topic {i+1} should be a dictionary")
                return False
            
            if "topic" not in topic:
                messagebox.showerror("Invalid JSON", f"{source_name}: Topic {i+1} missing 'topic' field")
                return False
                
            if "slides" not in topic or not isinstance(topic["slides"], list):
                messagebox.showerror("Invalid JSON", f"{source_name}: Topic {i+1} missing or invalid 'slides' field")
                return False
        
        return True

    def clear_all_files(self):
        """Clear all loaded files and topics"""
        if not self.all_topics_data:
            return
        
        result = messagebox.askyesno("Clear All", "Are you sure you want to clear all loaded content?")
        if result:
            self.all_topics_data = []
            self.loaded_files = []
            self.update_ui_after_load()
            self.status_label.config(text="🗑️ All content cleared", fg="orange")

    def update_ui_after_load(self):
        """Update UI elements after loading files"""
        # Update files listbox
        self.files_listbox.delete(0, tk.END)
        for source in self.loaded_files:
            self.files_listbox.insert(tk.END, source)
        
        # Update topics listbox and info
        self.topics_listbox.delete(0, tk.END)
        total_slides = 0
        total_code_slides = 0
        
        for i, topic in enumerate(self.all_topics_data, 1):
            topic_name = topic.get("topic", f"Topic {i}")
            slides = topic.get("slides", [])
            slides_count = len(slides)
            code_slides = sum(1 for slide in slides if slide.get("slide_type") == "code")
            source_file = topic.get("_source_file", "Unknown")
            total_slides += slides_count
            total_code_slides += code_slides
            
            display_text = f"{i}. {topic_name} ({slides_count} slides"
            if code_slides > 0:
                display_text += f", {code_slides} code"
            display_text += f") - {source_file}"
            
            self.topics_listbox.insert(tk.END, display_text)
        
        # Update info label
        if self.all_topics_data:
            source_count = len(self.loaded_files)
            info_text = f"📊 Total: {len(self.all_topics_data)} topics, {total_slides} slides"
            if total_code_slides > 0 and self.generate_code_files_var.get():
                info_text += f", {total_code_slides} code files will be generated"
            info_text += f" from {source_count} sources"
            self.topics_info_label.config(text=info_text, foreground="blue")
        else:
            self.topics_info_label.config(text="No topics loaded", foreground="gray")

    def browse_output_directory(self):
        """Browse for output directory"""
        directory = filedialog.askdirectory(title="Select Output Directory")
        if directory:
            self.output_dir_var.set(directory)

    def sanitize_filename(self, filename):
        """Sanitize filename for cross-platform compatibility"""
        # Remove or replace invalid characters
        invalid_chars = '<>:"/\\|?*'
        for char in invalid_chars:
            filename = filename.replace(char, '_')
        
        # Remove extra whitespace and limit length
        filename = re.sub(r'\s+', ' ', filename.strip())
        return filename[:100]  # Limit to 100 characters

    def detect_programming_language(self, slide_data):
        """Detect programming language from slide data"""
        # First check if language is explicitly specified
        if "language" in slide_data:
            return slide_data["language"].lower()
        
        # Try to detect from content
        content = slide_data.get("content", "")
        title = slide_data.get("title", "").lower()
        
        # Common language indicators in title
        title_indicators = {
            'python': ['python', 'py'], 'javascript': ['javascript', 'js', 'node'],
            'java': ['java'], 'cpp': ['c++', 'cpp'], 'c': ['c programming', ' c '],
            'csharp': ['c#', 'csharp'], 'php': ['php'], 'ruby': ['ruby'], 'go': ['golang', 'go'],
            'rust': ['rust'], 'swift': ['swift'], 'kotlin': ['kotlin'], 'sql': ['sql', 'database'],
            'html': ['html'], 'css': ['css'], 'typescript': ['typescript', 'ts'],
            'bash': ['bash', 'shell'], 'powershell': ['powershell'], 'yaml': ['yaml', 'yml'],
            'json': ['json'], 'xml': ['xml']
        }
        
        for lang, indicators in title_indicators.items():
            if any(indicator in title for indicator in indicators):
                return lang
        
        # Content-based detection
        content_lower = content.lower()
        
        # Check for specific syntax patterns
        if re.search(r'\bdef\s+\w+\s*\(.*\):', content) or 'import ' in content_lower or 'print(' in content_lower:
            return 'python'
        elif re.search(r'\bfunction\s+\w+\s*\(.*\)', content) or 'console.log' in content_lower or 'let ' in content_lower:
            return 'javascript'
        elif re.search(r'\bpublic\s+(static\s+)?void\s+main', content) or 'System.out.println' in content:
            return 'java'
        elif '#include' in content_lower or 'std::' in content_lower:
            return 'cpp'
        elif 'using System' in content or 'Console.WriteLine' in content:
            return 'csharp'
        elif '<?php' in content_lower or '$' in content and 'echo' in content_lower:
            return 'php'
        elif re.search(r'\bfunc\s+\w+\s*\(.*\)', content) and 'fmt.Println' in content:
            return 'go'
        elif 'fn ' in content_lower and 'println!' in content_lower:
            return 'rust'
        elif 'SELECT' in content.upper() or 'INSERT' in content.upper() or 'UPDATE' in content.upper():
            return 'sql'
        elif '<html>' in content_lower or '<!DOCTYPE' in content_lower:
            return 'html'
        elif re.search(r'\{[^}]*:[^}]*\}', content) and ('color:' in content_lower or 'font-' in content_lower):
            return 'css'
        elif content_lower.strip().startswith('{') and content_lower.strip().endswith('}'):
            return 'json'
        elif '#!/bin/bash' in content_lower or 'echo ' in content_lower:
            return 'bash'
        
        # Default fallback
        return 'txt'

    def get_file_extension(self, language):
        """Get file extension for a programming language"""
        return LANGUAGE_EXTENSIONS.get(language.lower(), '.txt')

    def create_code_file(self, folder_path, topic_name, slide_data, slide_index):
        """Create a separate code file for code slides"""
        if not self.generate_code_files_var.get():
            return None
        
        try:
            # Detect programming language
            language = self.detect_programming_language(slide_data)
            extension = self.get_file_extension(language)
            
            # Create filename
            slide_title = slide_data.get("title", f"slide_{slide_index}")
            sanitized_title = self.sanitize_filename(slide_title)
            
            # Special handling for Dockerfile
            if language == 'dockerfile':
                filename = "Dockerfile"
            else:
                filename = f"{slide_index}_{sanitized_title}{extension}"
            
            file_path = os.path.join(folder_path, filename)
            
            # Get code content
            content = slide_data.get("content", "")
            
            # Clean code content - remove markdown code blocks if present
            content_stripped = content.strip()
            # **FIXED HERE: Added closing quote**
            if content_stripped.startswith("``````"):
                lines = content_stripped.split('\n')
                # This handles ``````
                content = '\n'.join(lines[1:-1])
            
            # Write code file
            with open(file_path, "w", encoding="utf-8") as f:
                # Add header comment based on language
                if language in ['python', 'bash', 'shell', 'ruby', 'perl']:
                    f.write(f"# {slide_data.get('title', 'Code Example')}\n")
                    f.write(f"# Generated from: {topic_name}\n\n")
                elif language in ['javascript', 'java', 'cpp', 'c', 'csharp', 'php', 'go', 'rust', 'swift', 'kotlin', 'scala', 'css', 'typescript', 'dart', 'groovy']:
                    f.write(f"// {slide_data.get('title', 'Code Example')}\n")
                    f.write(f"// Generated from: {topic_name}\n\n")
                elif language in ['html', 'xml']:
                    f.write(f"<!-- {slide_data.get('title', 'Code Example')} -->\n")
                    f.write(f"<!-- Generated from: {topic_name} -->\n\n")
                elif language == 'sql':
                    f.write(f"-- {slide_data.get('title', 'Code Example')}\n")
                    f.write(f"-- Generated from: {topic_name}\n\n")
                
                # Write the actual code content
                f.write(content)
            
            return filename
            
        except Exception as e:
            print(f"Error creating code file: {str(e)}")
            return None

    def generate_all_presentations(self):
        """Generate PPTs for all loaded topics organized in folders with code files"""
        if not self.all_topics_data:
            messagebox.showwarning("Warning", "Please load at least one JSON file or add content from text")
            return
        
        try:
            starting_number = int(self.folder_number_var.get())
        except ValueError:
            messagebox.showerror("Error", "Please enter a valid starting folder number")
            return
        
        output_base_dir = self.output_dir_var.get()
        if not os.path.exists(output_base_dir):
            messagebox.showerror("Error", "Output directory does not exist")
            return
        
        theme = THEMES.get(self.theme_var.get(), THEMES["Light"])
        naming_style = self.naming_var.get()
        
        try:
            created_folders = []
            total_topics = len(self.all_topics_data)
            total_code_files = 0
            
            for i, topic_data in enumerate(self.all_topics_data):
                folder_number = starting_number + i
                topic_name = topic_data.get("topic", f"Topic {i+1}")
                sanitized_topic = self.sanitize_filename(topic_name)
                
                # Create folder name based on naming style
                if naming_style == "numbered":
                    folder_name = f"{folder_number:02d}. {sanitized_topic}"
                else:
                    folder_name = sanitized_topic
                
                folder_path = os.path.join(output_base_dir, folder_name)
                os.makedirs(folder_path, exist_ok=True)
                created_folders.append(folder_name)
                
                # Create code files for code slides
                code_files_created = []
                slides = topic_data.get("slides", [])
                
                for slide_index, slide_data in enumerate(slides, 1):
                    if slide_data.get("slide_type") == "code":
                        code_filename = self.create_code_file(folder_path, topic_name, slide_data, slide_index)
                        if code_filename:
                            code_files_created.append(code_filename)
                            total_code_files += 1
                
                # Create PPT
                ppt_filename = f"{sanitized_topic}.pptx"
                ppt_path = os.path.join(folder_path, ppt_filename)
                
                # Generate presentation for this topic
                self.create_single_presentation(topic_data, ppt_path, theme)
                
                # Update status
                progress_text = f"Processing... {i+1}/{total_topics} - {topic_name[:30]}{'...' if len(topic_name) > 30 else ''}"
                if code_files_created:
                    progress_text += f" ({len(code_files_created)} code files)"
                self.status_label.config(text=progress_text, fg="blue")
                self.root.update()
            
            # Create a summary file
            self.create_summary_file(output_base_dir, created_folders, total_code_files)
            
            # Success message
            success_msg = f"🎉 Successfully created {len(created_folders)} PPTs"
            if total_code_files > 0:
                success_msg += f" + {total_code_files} code files"
            success_msg += f"!\n\n📂 Location: {output_base_dir}\n"
            success_msg += f"📊 From {len(self.loaded_files)} sources\n\n"
            success_msg += "📁 Created folders:\n"
            success_msg += "\n".join([f"-  {folder}" for folder in created_folders[:8]])
            if len(created_folders) > 8:
                success_msg += f"\n... and {len(created_folders) - 8} more folders"
            
            messagebox.showinfo("Success! 🎉", success_msg)
            
            status_text = f"✅ Created {len(created_folders)} PPTs from {len(self.loaded_files)} sources"
            if total_code_files > 0:
                status_text += f" + {total_code_files} code files"
            self.status_label.config(text=status_text, fg="green")
            
        except Exception as e:
            messagebox.showerror("Error", f"An error occurred during generation:\n{str(e)}")
            self.status_label.config(text="❌ Error occurred during generation", fg="red")

    def create_summary_file(self, output_dir, folders, total_code_files):
        """Create a summary text file listing all generated presentations and code files"""
        summary_path = os.path.join(output_dir, "Generated_PPTs_Summary.txt")
        
        with open(summary_path, "w", encoding="utf-8") as f:
            f.write("📱 Mobile PPT Generator - Summary Report\n")
            f.write("=" * 50 + "\n\n")
            f.write(f"📅 Generated on: {datetime.datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n")
            f.write(f"📂 Output Directory: {output_dir}\n")
            f.write(f"🎨 Theme Used: {self.theme_var.get()}\n")
            f.write(f"📊 Total Topics: {len(self.all_topics_data)}\n")
            f.write(f"📄 Source Count: {len(self.loaded_files)}\n")
            if total_code_files > 0:
                f.write(f"💻 Total Code Files: {total_code_files}\n")
            f.write(f"🔧 Code File Generation: {'Enabled' if self.generate_code_files_var.get() else 'Disabled'}\n\n")
            
            f.write("📁 Generated Folders and Content:\n")
            f.write("-" * 40 + "\n")
            
            for i, (folder, topic_data) in enumerate(zip(folders, self.all_topics_data), 1):
                topic_name = topic_data.get("topic", f"Topic {i}")
                slides = topic_data.get("slides", [])
                slides_count = len(slides)
                code_slides = [s for s in slides if s.get("slide_type") == "code"]
                source_file = topic_data.get("_source_file", "Unknown")
                
                f.write(f"{i:2d}. {folder}/\n")
                f.write(f"    📄 {self.sanitize_filename(topic_name)}.pptx ({slides_count} slides)\n")
                
                # List code files if any
                if code_slides and self.generate_code_files_var.get():
                    f.write(f"    💻 Code files ({len(code_slides)}):\n")
                    slide_num_map = {id(slide): idx + 1 for idx, slide in enumerate(slides)}
                    for slide_index, slide in enumerate(code_slides, 1):
                        language = self.detect_programming_language(slide)
                        extension = self.get_file_extension(language)
                        slide_title = slide.get("title", f"slide_{slide_index}")
                        sanitized_title = self.sanitize_filename(slide_title)
                        if language == 'dockerfile':
                            filename = "Dockerfile"
                        else:
                            original_index = slide_num_map.get(id(slide), slide_index)
                            filename = f"{original_index}_{sanitized_title}{extension}"
                        f.write(f"        -  {filename} ({language})\n")
                
                f.write(f"    📋 Source: {source_file}\n\n")
            
            f.write("\n🚀 All presentations are mobile-optimized with:\n")
            f.write("-  Portrait orientation (9:16 aspect ratio)\n")
            f.write("-  Large fonts for mobile readability\n")
            f.write("-  Proper word wrapping and spacing\n")
            f.write("-  Code syntax highlighting\n")
            f.write("-  Summary slides for each topic\n")
            if self.generate_code_files_var.get():
                f.write("-  Separate code files for all 'slide_type': 'code' slides\n")
                f.write("-  Automatic programming language detection\n")
                f.write("-  Proper file extensions based on detected language\n")

    def create_single_presentation(self, topic_data, output_path, theme):
        """Create a single presentation for a topic"""
        # Create new presentation with mobile dimensions
        prs = Presentation()
        # Mobile-friendly aspect ratio (9:16 portrait)
        prs.slide_width = Cm(20.32)  # 8 inches
        prs.slide_height = Cm(36.07)  # 14.2 inches

        topic_name = topic_data.get("topic", "Untitled Topic")
        slides_data = topic_data.get("slides", [])
        
        # Add title slide
        self.add_title_slide(prs, topic_name)
        
        # Process all slides
        slide_titles = []
        for slide_data in slides_data:
            title = slide_data.get("title", "Untitled")
            content = slide_data.get("content", "")
            slide_type = slide_data.get("slide_type", "text")

            slide_titles.append(title)

            if slide_type == "code":
                self.add_code_slide(prs, title, content, theme)
            elif slide_type == "table" and isinstance(content, list):
                self.add_table_slide(prs, title, content, theme)
            else:
                self.add_text_slide(prs, title, content, theme)

        # Add summary slide
        self.add_summary_slide(prs, slide_titles, theme)

        # Save presentation
        prs.save(output_path)

    def add_title_slide(self, prs, title):
        """Add mobile-optimized title slide"""
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Main title - larger and centered for mobile
        title_box = slide.shapes.add_textbox(Cm(2), Cm(8), Cm(16.32), Cm(8))
        title_frame = title_box.text_frame
        title_frame.clear()
        title_frame.word_wrap = True
        
        p = title_frame.add_paragraph()
        p.text = title
        p.font.size = Pt(42)  # Large font for mobile
        p.font.bold = True
        p.font.name = "Arial"
        p.alignment = PP_ALIGN.CENTER
        
        # Subtitle
        subtitle_box = slide.shapes.add_textbox(Cm(2), Cm(18), Cm(16.32), Cm(6))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.clear()
        
        p2 = subtitle_frame.add_paragraph()
        p2.text = "📱 Mobile-Optimized Presentation"
        p2.font.size = Pt(26)
        p2.font.name = "Arial"
        p2.alignment = PP_ALIGN.CENTER
        
        # Footer with generator info
        footer_box = slide.shapes.add_textbox(Cm(2), Cm(32), Cm(16.32), Cm(3))
        footer_frame = footer_box.text_frame
        footer_frame.clear()
        
        p3 = footer_frame.add_paragraph()
        p3.text = "Generated by Final Mobile PPT Generator"
        p3.font.size = Pt(16)
        p3.font.name = "Arial"
        p3.alignment = PP_ALIGN.CENTER
        p3.font.color.rgb = RGBColor(128, 128, 128)

    def add_text_slide(self, prs, title, content, theme):
        """Add mobile-optimized text slide with proper word wrapping"""
        max_words_per_slide = 150  # Reduced for mobile readability
        
        # Convert markdown to plain text
        if content:
            html = markdown2.markdown(content)
            soup = BeautifulSoup(html, "html.parser")
            words = soup.get_text().split()
        else:
            words = ["No content available"]

        # Split content across multiple slides if needed
        for i in range(0, len(words), max_words_per_slide):
            slide_layout = prs.slide_layouts[6] # Blank layout
            slide = prs.slides.add_slide(slide_layout)
            
            # Add title
            slide_title = title if i == 0 else f"{title} (Part {(i//max_words_per_slide) + 1})"
            self.add_slide_title(slide, slide_title, theme)
            
            # Add content text box
            content_text = " ".join(words[i:i + max_words_per_slide])
            textbox = slide.shapes.add_textbox(Cm(2), Cm(8), Cm(16.32), Cm(25))
            tf = textbox.text_frame
            tf.clear()
            tf.word_wrap = True
            
            p = tf.add_paragraph()
            p.text = content_text
            p.font.name = "Arial"
            p.font.size = Pt(28)  # Large font for mobile readability
            p.font.color.rgb = RGBColor(*self.hex_to_rgb(theme["text_color"]))
            p.alignment = PP_ALIGN.LEFT
            p.line_spacing = 1.2  # Better line spacing for mobile

    def add_code_slide(self, prs, title, code, theme):
        """Add mobile-optimized code slide"""
        max_lines = 20  # Increased for mobile vertical space
        
        # Clean code content - remove markdown code blocks
        code_content = code.strip() if code else "# No code content"
        # **FIXED HERE: Added closing quote**
        if code_content.startswith("``````"):
            lines = code_content.split('\n')
            code_content = '\n'.join(lines[1:-1])
        
        lines = code_content.split("\n")
        
        for i in range(0, len(lines), max_lines):
            slide_layout = prs.slide_layouts[6] # Blank layout
            slide = prs.slides.add_slide(slide_layout)
            
            # Add title
            slide_title = title if i == 0 else f"{title} (Part {(i//max_lines) + 1})"
            self.add_slide_title(slide, slide_title, theme)
            
            # Code content
            chunk = "\n".join(lines[i:i + max_lines])
            font_size = self.calculate_mobile_code_font_size(chunk)
            
            # Code text box with background
            textbox = slide.shapes.add_textbox(Cm(1.5), Cm(8), Cm(17.32), Cm(25))
            
            # Set background color
            textbox.fill.solid()
            if theme["bg_color"] == "#FFFFFF":
                textbox.fill.fore_color.rgb = RGBColor(245, 245, 245)  # Light gray for light theme
            else:
                textbox.fill.fore_color.rgb = RGBColor(30, 30, 30) # Dark gray for dark themes
            
            tf = textbox.text_frame
            tf.clear()
            tf.word_wrap = False  # Preserve code formatting
            tf.margin_left = Cm(0.5)
            tf.margin_right = Cm(0.5)
            tf.margin_top = Cm(0.5)
            tf.margin_bottom = Cm(0.5)
            
            p = tf.add_paragraph()
            p.text = chunk
            p.font.name = "Consolas"  # Better mobile code font
            p.font.size = Pt(font_size)
            p.font.color.rgb = RGBColor(*self.hex_to_rgb(theme["text_color"]))
            p.alignment = PP_ALIGN.LEFT

    def add_table_slide(self, prs, title, table_data, theme):
        """Add mobile-optimized table slide"""
        if not table_data or not isinstance(table_data, list) or len(table_data) == 0 or not all(isinstance(row, dict) for row in table_data):
            self.add_text_slide(prs, title, "Invalid table data format. Expected a list of dictionaries.", theme)
            return
            
        slide_layout = prs.slide_layouts[6] # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Add title
        self.add_slide_title(slide, title, theme)
        
        # Table dimensions
        headers = list(table_data[0].keys())
        rows = len(table_data)
        cols = len(headers)
        
        # Create table with mobile-friendly sizing
        table_height = min(Cm(25), Cm(3 + (rows * 2)))  # Dynamic height with max limit
        table_shape = slide.shapes.add_table(
            rows + 1, cols, 
            Cm(1), Cm(9), 
            Cm(18.32), table_height
        )
        table = table_shape.table
        
        # Header row
        for col_idx, header in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = str(header)
            
            # Header styling
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.bold = True
            paragraph.font.size = Pt(20)  # Larger font for mobile
            paragraph.font.name = "Arial"
            paragraph.font.color.rgb = RGBColor(255, 255, 255) # White text for contrast
            paragraph.alignment = PP_ALIGN.CENTER
            
            # Header background
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(*self.hex_to_rgb(theme["accent"]))
        
        # Data rows
        for row_idx, row_data in enumerate(table_data):
            for col_idx, header in enumerate(headers):
                cell_value = row_data.get(header, "")
                cell = table.cell(row_idx + 1, col_idx)
                cell.text = str(cell_value)
                
                # Data cell styling
                paragraph = cell.text_frame.paragraphs[0]
                paragraph.font.size = Pt(18)  # Mobile-friendly size
                paragraph.font.name = "Arial"
                paragraph.font.color.rgb = RGBColor(*self.hex_to_rgb(theme["text_color"]))
                paragraph.alignment = PP_ALIGN.CENTER

    def add_summary_slide(self, prs, titles, theme):
        """Add mobile-optimized summary slide"""
        slide_layout = prs.slide_layouts[6] # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Add title
        self.add_slide_title(slide, "📋 Summary", theme)
        
        # Summary content
        textbox = slide.shapes.add_textbox(Cm(2), Cm(8), Cm(16.32), Cm(25))
        tf = textbox.text_frame
        tf.clear()
        tf.word_wrap = True
        
        for i, title in enumerate(titles, 1):
            p = tf.add_paragraph()
            p.text = f"{i}. {title}"
            p.font.name = "Arial"
            p.font.size = Pt(24)  # Large font for mobile
            p.font.color.rgb = RGBColor(*self.hex_to_rgb(theme["text_color"]))
            p.level = 0
            p.space_after = Pt(12)  # Spacing between items

    def add_slide_title(self, slide, title, theme):
        """Add consistent mobile-optimized title to slide"""
        # Create custom title box
        title_box = slide.shapes.add_textbox(Cm(1), Cm(2), Cm(18.32), Cm(4))
        title_frame = title_box.text_frame
        title_frame.clear()
        title_frame.word_wrap = True
        
        p = title_frame.add_paragraph()
        p.text = title
        p.font.size = Pt(32)  # Large mobile-friendly title
        p.font.bold = True
        p.font.name = "Arial"
        p.font.color.rgb = RGBColor(*self.hex_to_rgb(theme["accent"]))
        p.alignment = PP_ALIGN.CENTER

    def calculate_mobile_code_font_size(self, code):
        """Calculate appropriate font size for mobile code display"""
        lines = code.count('\n') + 1
        max_line_length = max(len(line) for line in code.split('\n')) if code else 0
        
        # Adjust based on both line count and line length for mobile
        if lines <= 10 and max_line_length <= 50:
            return 22
        elif lines <= 15 and max_line_length <= 60:
            return 20
        elif lines <= 20 and max_line_length <= 70:
            return 18
        else:
            return 16

    def hex_to_rgb(self, hex_color):
        """Convert hex color to RGB tuple"""
        hex_color = hex_color.lstrip("#")
        return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))


if __name__ == "__main__":
    root = tk.Tk()
    app = FinalPPTCreatorApp(root)
    root.mainloop()
