import os
import json
import tkinter as tk
from tkinter import filedialog, ttk, messagebox
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN, PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor
import pygments
from pygments import lexers, formatters
from bs4 import BeautifulSoup
import markdown2

# Mobile-optimized themes with better contrast
THEMES = {
    "Light": {"bg_color": "#FFFFFF", "text_color": "#000000", "accent": "#007ACC"},
    "Dark": {"bg_color": "#1E1E1E", "text_color": "#FFFFFF", "accent": "#4FC3F7"},
    "Blue": {"bg_color": "#0D47A1", "text_color": "#FFFFFF", "accent": "#81D4FA"}
}

class PPTCreatorApp:
    def __init__(self, root):
        self.root = root
        self.root.title("📱 Mobile PPT Generator")
        self.root.geometry("400x300")
        self.prs = Presentation()
        
        # Setup UI
        self.setup_ui()
        
    def setup_ui(self):
        # Main frame
        main_frame = ttk.Frame(self.root, padding="20")
        main_frame.pack(fill=tk.BOTH, expand=True)
        
        # Title
        title_label = ttk.Label(main_frame, text="📱 Mobile PPT Generator", 
                               font=("Arial", 16, "bold"))
        title_label.pack(pady=(0, 20))
        
        # Theme selection
        theme_frame = ttk.Frame(main_frame)
        theme_frame.pack(fill=tk.X, pady=(0, 20))
        
        ttk.Label(theme_frame, text="Select Theme:", font=("Arial", 12)).pack(anchor=tk.W)
        self.theme_var = tk.StringVar(value="Light")
        self.theme_dropdown = ttk.Combobox(theme_frame, textvariable=self.theme_var, 
                                          values=list(THEMES.keys()), state="readonly")
        self.theme_dropdown.pack(fill=tk.X, pady=(5, 0))
        
        # Load button
        load_btn = tk.Button(main_frame, text="📂 Load JSON & Create PPT", 
                           font=("Arial", 12), bg="#007ACC", fg="white",
                           command=self.load_json, cursor="hand2")
        load_btn.pack(pady=20, fill=tk.X)
        
        # Status label
        self.status_label = tk.Label(main_frame, text="Ready to create mobile-friendly presentations", 
                                   font=("Arial", 10), fg="gray")
        self.status_label.pack(pady=(10, 0))

    def load_json(self):
        file_paths = filedialog.askopenfilenames(
            filetypes=[("JSON Files", "*.json")], 
            title="Select JSON Files"
        )
        if not file_paths:
            return

        try:
            for file_path in file_paths:
                self.process_json_file(file_path)
        except Exception as e:
            messagebox.showerror("Error", f"An error occurred: {str(e)}")
            self.status_label.config(text="❌ Error occurred during processing", fg="red")

    def process_json_file(self, file_path):
        with open(file_path, "r", encoding="utf-8") as file:
            slides_data = json.load(file)

        if not slides_data:
            messagebox.showwarning("Warning", "No slide data found in JSON file")
            return

        # Create new presentation with mobile dimensions
        self.prs = Presentation()
        # Mobile-friendly aspect ratio (9:16 portrait)
        self.prs.slide_width = Cm(20.32)  # 8 inches
        self.prs.slide_height = Cm(36.07)  # 14.2 inches

        # Get selected theme
        theme = THEMES.get(self.theme_var.get(), THEMES["Light"])
        
        # Add title slide
        if slides_data:
            self.add_title_slide(slides_data[0].get("title", "Mobile Presentation"))

        # Process all slides
        slide_titles = []
        for slide_data in slides_data:
            title = slide_data.get("title", "Untitled")
            content = slide_data.get("content", "")
            slide_type = slide_data.get("slide_type", "text")

            slide_titles.append(title)

            if slide_type == "code":
                self.add_code_slide(title, content, theme)
            elif slide_type == "table" and isinstance(content, list):
                self.add_table_slide(title, content, theme)
            else:
                self.add_text_slide(title, content, theme)

        # Add summary slide
        self.add_summary_slide(slide_titles, theme)

        # Save presentation
        self.save_presentation(file_path)

    def add_title_slide(self, title):
        """Add mobile-optimized title slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[0])
        
        # Clear existing placeholders and create custom layout
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                shape.text = ""
        
        # Main title - larger and centered for mobile
        title_box = slide.shapes.add_textbox(Cm(2), Cm(8), Cm(16.32), Cm(6))
        title_frame = title_box.text_frame
        title_frame.clear()
        title_frame.word_wrap = True
        
        p = title_frame.add_paragraph()
        p.text = title
        p.font.size = Pt(40)  # Large font for mobile
        p.font.bold = True
        p.font.name = "Arial"
        p.alignment = PP_ALIGN.CENTER
        
        # Subtitle
        subtitle_box = slide.shapes.add_textbox(Cm(2), Cm(16), Cm(16.32), Cm(4))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.clear()
        
        p2 = subtitle_frame.add_paragraph()
        p2.text = "📱 Mobile-Optimized Presentation"
        p2.font.size = Pt(24)
        p2.font.name = "Arial"
        p2.alignment = PP_ALIGN.CENTER

    def add_text_slide(self, title, content, theme):
        """Add mobile-optimized text slide with proper word wrapping"""
        max_words_per_slide = 150  # Reduced for mobile readability
        
        # Convert markdown to plain text
        if content:
            html = markdown2.markdown(content)
            soup = BeautifulSoup(html, "html.parser")
            words = soup.get_text().split()
        else:
            words = ["No content available"]

        # Split content across multiple slides if needed
        for i in range(0, len(words), max_words_per_slide):
            slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
            
            # Add title
            slide_title = title if i == 0 else f"{title} (Part {(i//max_words_per_slide) + 1})"
            self.add_slide_title(slide, slide_title, theme)
            
            # Add content text box
            content_text = " ".join(words[i:i + max_words_per_slide])
            textbox = slide.shapes.add_textbox(Cm(2), Cm(8), Cm(16.32), Cm(25))
            tf = textbox.text_frame
            tf.clear()
            tf.word_wrap = True
            tf.margin_left = Cm(1)
            tf.margin_right = Cm(1)
            tf.margin_top = Cm(1)
            tf.margin_bottom = Cm(1)
            
            p = tf.add_paragraph()
            p.text = content_text
            p.font.name = "Arial"
            p.font.size = Pt(28)  # Large font for mobile readability
            p.font.color.rgb = RGBColor(*self.hex_to_rgb(theme["text_color"]))
            p.alignment = PP_ALIGN.LEFT
            p.line_spacing = 1.2  # Better line spacing for mobile

    def add_code_slide(self, title, code, theme):
        """Add mobile-optimized code slide"""
        max_lines = 20  # Increased for mobile vertical space
        lines = code.split("\n") if code else ["# No code content"]
        
        for i in range(0, len(lines), max_lines):
            slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
            
            # Add title
            slide_title = title if i == 0 else f"{title} (Part {(i//max_lines) + 1})"
            self.add_slide_title(slide, slide_title, theme)
            
            # Code content
            chunk = "\n".join(lines[i:i + max_lines])
            font_size = self.calculate_mobile_code_font_size(chunk)
            
            # Code text box with background
            textbox = slide.shapes.add_textbox(Cm(1.5), Cm(8), Cm(17.32), Cm(25))
            
            # Set background color
            textbox.fill.solid()
            if theme["bg_color"] == "#FFFFFF":
                textbox.fill.fore_color.rgb = RGBColor(245, 245, 245)  # Light gray for light theme
            else:
                textbox.fill.fore_color.rgb = RGBColor(*self.hex_to_rgb(theme["bg_color"]))
            
            tf = textbox.text_frame
            tf.clear()
            tf.word_wrap = False  # Preserve code formatting
            tf.margin_left = Cm(0.5)
            tf.margin_right = Cm(0.5)
            tf.margin_top = Cm(0.5)
            tf.margin_bottom = Cm(0.5)
            
            p = tf.add_paragraph()
            p.text = chunk
            p.font.name = "Consolas"  # Better mobile code font
            p.font.size = Pt(font_size)
            p.font.color.rgb = RGBColor(*self.hex_to_rgb(theme["text_color"]))
            p.alignment = PP_ALIGN.LEFT

    def add_table_slide(self, title, table_data, theme):
        """Add mobile-optimized table slide"""
        if not table_data or not isinstance(table_data, list):
            return
            
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        
        # Add title
        self.add_slide_title(slide, title, theme)
        
        # Table dimensions
        rows = len(table_data)
        if rows == 0:
            return
            
        cols = len(table_data[0]) if isinstance(table_data[0], dict) else 0
        if cols == 0:
            return
        
        # Create table with mobile-friendly sizing
        table_height = min(Cm(25), Cm(3 + (rows * 2)))  # Dynamic height with max limit
        table = slide.shapes.add_table(
            rows + 1, cols, 
            Cm(1), Cm(9), 
            Cm(18.32), table_height
        ).table
        
        # Header row
        headers = list(table_data[0].keys())
        for col_idx, header in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = str(header)
            
            # Header styling
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.bold = True
            paragraph.font.size = Pt(20)  # Larger font for mobile
            paragraph.font.name = "Arial"
            paragraph.font.color.rgb = RGBColor(*self.hex_to_rgb(theme["text_color"]))
            paragraph.alignment = PP_ALIGN.CENTER
            
            # Header background
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(*self.hex_to_rgb(theme["accent"]))
        
        # Data rows
        for row_idx, row_data in enumerate(table_data):
            for col_idx, value in enumerate(row_data.values()):
                cell = table.cell(row_idx + 1, col_idx)
                cell.text = str(value)
                
                # Data cell styling
                paragraph = cell.text_frame.paragraphs[0]
                paragraph.font.size = Pt(18)  # Mobile-friendly size
                paragraph.font.name = "Arial"
                paragraph.font.color.rgb = RGBColor(*self.hex_to_rgb(theme["text_color"]))
                paragraph.alignment = PP_ALIGN.CENTER

    def add_summary_slide(self, titles, theme):
        """Add mobile-optimized summary slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        
        # Add title
        self.add_slide_title(slide, "📋 Summary", theme)
        
        # Summary content
        textbox = slide.shapes.add_textbox(Cm(2), Cm(8), Cm(16.32), Cm(25))
        tf = textbox.text_frame
        tf.clear()
        tf.word_wrap = True
        
        for i, title in enumerate(titles, 1):
            p = tf.add_paragraph()
            p.text = f"{i}. {title}"
            p.font.name = "Arial"
            p.font.size = Pt(24)  # Large font for mobile
            p.font.color.rgb = RGBColor(*self.hex_to_rgb(theme["text_color"]))
            p.level = 0
            p.space_after = Pt(12)  # Spacing between items

    def add_slide_title(self, slide, title, theme):
        """Add consistent mobile-optimized title to slide"""
        # Clear existing title placeholders
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame') and shape.text_frame:
                shape.text_frame.clear()
        
        # Create custom title box
        title_box = slide.shapes.add_textbox(Cm(1), Cm(2), Cm(18.32), Cm(4))
        title_frame = title_box.text_frame
        title_frame.clear()
        title_frame.word_wrap = True
        
        p = title_frame.add_paragraph()
        p.text = title
        p.font.size = Pt(32)  # Large mobile-friendly title
        p.font.bold = True
        p.font.name = "Arial"
        p.font.color.rgb = RGBColor(*self.hex_to_rgb(theme["accent"]))
        p.alignment = PP_ALIGN.CENTER

    def calculate_mobile_code_font_size(self, code):
        """Calculate appropriate font size for mobile code display"""
        lines = code.count('\n') + 1
        max_line_length = max(len(line) for line in code.split('\n')) if code else 0
        
        # Adjust based on both line count and line length for mobile
        if lines <= 10 and max_line_length <= 50:
            return 22
        elif lines <= 15 and max_line_length <= 60:
            return 20
        elif lines <= 20 and max_line_length <= 70:
            return 18
        else:
            return 16

    def hex_to_rgb(self, hex_color):
        """Convert hex color to RGB tuple"""
        hex_color = hex_color.lstrip("#")
        return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

    def save_presentation(self, original_file_path):
        """Save the presentation with user-selected filename"""
        default_name = f"mobile_ppt_{os.path.splitext(os.path.basename(original_file_path))[0]}.pptx"
        save_path = filedialog.asksaveasfilename(
            defaultextension=".pptx",
            filetypes=[("PowerPoint Files", "*.pptx")],
            title="Save Mobile PPT"
        )
        
        if save_path:
            self.prs.save(save_path)
            self.status_label.config(
                text=f"✅ Mobile PPT saved: {os.path.basename(save_path)}", 
                fg="green"
            )
            messagebox.showinfo("Success", f"Mobile-friendly presentation saved successfully!\n\nFile: {save_path}")
        else:
            self.status_label.config(text="❌ Save cancelled", fg="orange")

if __name__ == "__main__":
    root = tk.Tk()
    app = PPTCreatorApp(root)
    root.mainloop()